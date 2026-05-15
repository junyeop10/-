"""Office document text extraction helpers."""

from __future__ import annotations

from pathlib import Path

from src.text_cleaner import build_sampled_text


def extract_docx_text(
    path: str | Path,
    total_limit: int = 4500,
    part_limit: int = 1500,
) -> str:
    """Extract paragraph and table text from a DOCX file."""
    from docx import Document

    document = Document(str(path))
    chunks: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            chunks.append(text)

    for table in document.tables:
        for row in table.rows:
            row_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_values:
                chunks.append(" | ".join(row_values))

    return build_sampled_text("\n".join(chunks), total_limit=total_limit, part_limit=part_limit)


def extract_xlsx_text(
    path: str | Path,
    fast: bool = True,
    total_limit: int = 4500,
    part_limit: int = 1500,
) -> str:
    """Extract sheet names and top rows from an XLSX file."""
    from openpyxl import load_workbook

    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    chunks: list[str] = []

    max_rows = 20 if fast else 100
    max_cols = 8

    for sheet in workbook.worksheets:
        chunks.append(f"sheet: {sheet.title}")
        for row in sheet.iter_rows(min_row=1, max_row=max_rows, values_only=True):
            row_values = [str(value).strip() for value in row[:max_cols] if value not in (None, "")]
            if row_values:
                chunks.append(" | ".join(row_values))

    return build_sampled_text("\n".join(chunks), total_limit=total_limit, part_limit=part_limit)


def extract_pptx_text(
    path: str | Path,
    total_limit: int = 4500,
    part_limit: int = 1500,
) -> str:
    """Extract visible text from PPTX slides."""
    from pptx import Presentation

    presentation = Presentation(str(path))
    chunks: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        chunks.append(f"slide: {slide_index}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text).strip()
                if text:
                    chunks.append(text)

    return build_sampled_text("\n".join(chunks), total_limit=total_limit, part_limit=part_limit)
