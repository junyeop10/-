from __future__ import annotations

import unittest
import zipfile
from pathlib import Path
from shutil import rmtree
from uuid import uuid4

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from src.file_reader import SUPPORTED_SUFFIXES, discover_supported_files, extract_text_from_file


class FileReaderTest(unittest.TestCase):
    def setUp(self) -> None:
        self.base_dir = Path("tests_runtime") / f"reader_case_{uuid4().hex}"
        self.base_dir.mkdir(parents=True, exist_ok=True)

    def tearDown(self) -> None:
        if self.base_dir.exists():
            rmtree(self.base_dir, ignore_errors=True)

    def test_supported_suffixes_include_office_formats(self) -> None:
        self.assertIn(".docx", SUPPORTED_SUFFIXES)
        self.assertIn(".xlsx", SUPPORTED_SUFFIXES)
        self.assertIn(".pptx", SUPPORTED_SUFFIXES)
        self.assertIn(".hwp", SUPPORTED_SUFFIXES)
        self.assertIn(".hwpx", SUPPORTED_SUFFIXES)

    def test_discover_supported_files_finds_office_documents(self) -> None:
        (self.base_dir / "sample.txt").write_text("report summary", encoding="utf-8")
        self._create_docx(self.base_dir / "sample.docx")
        self._create_xlsx(self.base_dir / "sample.xlsx")
        self._create_pptx(self.base_dir / "sample.pptx")
        self._create_hwpx(self.base_dir / "sample.hwpx")

        files = discover_supported_files(self.base_dir)
        suffixes = sorted(path.suffix.lower() for path in files)
        self.assertEqual(suffixes, [".docx", ".hwpx", ".pptx", ".txt", ".xlsx"])

    def test_extract_text_from_docx(self) -> None:
        path = self.base_dir / "contract.docx"
        self._create_docx(path)

        text = extract_text_from_file(path, fast=True)
        self.assertIn("contract agreement", text.lower())
        self.assertIn("payment term", text.lower())

    def test_extract_text_from_xlsx(self) -> None:
        path = self.base_dir / "invoice.xlsx"
        self._create_xlsx(path)

        text = extract_text_from_file(path, fast=True)
        self.assertIn("sheet: invoice", text.lower())
        self.assertIn("amount due", text.lower())

    def test_extract_text_from_pptx(self) -> None:
        path = self.base_dir / "presentation.pptx"
        self._create_pptx(path)

        text = extract_text_from_file(path, fast=True)
        self.assertIn("slide: 1", text.lower())
        self.assertIn("business plan presentation", text.lower())

    def test_extract_text_from_hwpx(self) -> None:
        path = self.base_dir / "meeting.hwpx"
        self._create_hwpx(path)

        text = extract_text_from_file(path, fast=True)
        self.assertIn("회의록", text)
        self.assertIn("결정사항", text)

    def _create_docx(self, path: Path) -> None:
        document = Document()
        document.add_paragraph("Contract agreement for project delivery")
        document.add_paragraph("Payment term and report summary")
        document.save(path)

    def _create_xlsx(self, path: Path) -> None:
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Invoice"
        sheet.append(["Invoice", "Amount Due", "Status"])
        sheet.append(["INV-001", 1000, "Open"])
        workbook.save(path)

    def _create_pptx(self, path: Path) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Business Plan Presentation"
        slide.placeholders[1].text = "Target market and execution plan"
        presentation.save(path)

    def _create_hwpx(self, path: Path) -> None:
        section_xml = """<?xml version="1.0" encoding="UTF-8"?>
<hs:sec xmlns:hs="http://www.hancom.co.kr/hwpml/2011/section"
        xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph">
  <hp:p><hp:run><hp:t>회의록</hp:t></hp:run></hp:p>
  <hp:p><hp:run><hp:t>안건 검토 및 결정사항 정리</hp:t></hp:run></hp:p>
</hs:sec>
"""
        with zipfile.ZipFile(path, "w") as archive:
            archive.writestr("Contents/section0.xml", section_xml)


if __name__ == "__main__":
    unittest.main()
