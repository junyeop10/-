"""
stage0_extract.py — Stage 0: 본문 텍스트 추출

[역할] PDF·DOCX·XLSX·PPTX 등에서 텍스트를 뽑아 앞/중/뒤 3분할 청크로 나눕니다.
[입력] file_bytes, filename, ext
[출력] {"status": "success"|"failed", "front", "middle", "rear", "method"}
[담당] 김준엽 (feature/stage1-extract)
[다음] stage2_ocr (실패 시) → stage4_embedding
"""

from io import BytesIO
from pathlib import Path

CHUNK_SIZE = 1500


def _split_text(text: str) -> tuple[str, str, str]:
    if not text:
        return "", "", ""
    length = len(text)
    third = max(length // 3, 1)
    part1 = text[:third]
    part2 = text[third : third * 2]
    part3 = text[third * 2 :]
    return (
        part1[:CHUNK_SIZE],
        part2[:CHUNK_SIZE],
        part3[:CHUNK_SIZE],
    )


def _extract_pdf(file_bytes: bytes) -> tuple[str, str]:
    import fitz

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    parts = []
    for page in doc:
        parts.append(page.get_text())
    doc.close()
    return "".join(parts), "pymupdf"


def _extract_docx(file_bytes: bytes) -> tuple[str, str]:
    from docx import Document

    doc = Document(BytesIO(file_bytes))
    text = "\n".join(p.text for p in doc.paragraphs if p.text)
    return text, "python-docx"


def _extract_xlsx(file_bytes: bytes) -> tuple[str, str]:
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(sheet.title)
        rows = list(sheet.iter_rows(max_row=1, values_only=True))
        if rows:
            header = " ".join(str(c) for c in rows[0] if c is not None)
            if header:
                parts.append(header)
    wb.close()
    return "\n".join(parts), "openpyxl"


def _extract_pptx(file_bytes: bytes) -> tuple[str, str]:
    from pptx import Presentation

    prs = Presentation(BytesIO(file_bytes))
    parts = []
    for slide in prs.slides:
        if slide.shapes.title and slide.shapes.title.text:
            parts.append(slide.shapes.title.text)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts), "python-pptx"


def run(file_bytes: bytes, filename: str, ext: str) -> dict:
    ext = ext.lower() if ext else Path(filename).suffix.lower()

    if ext in (".hwp", ".hwpx", ".jpg", ".jpeg", ".png"):
        return {
            "status": "failed",
            "front": "",
            "middle": "",
            "rear": "",
            "method": ext.lstrip("."),
        }

    try:
        if ext == ".pdf":
            text, method = _extract_pdf(file_bytes)
        elif ext == ".docx":
            text, method = _extract_docx(file_bytes)
        elif ext == ".xlsx":
            text, method = _extract_xlsx(file_bytes)
        elif ext == ".pptx":
            text, method = _extract_pptx(file_bytes)
        else:
            return {
                "status": "failed",
                "front": "",
                "middle": "",
                "rear": "",
                "method": "unknown",
            }

        if not text.strip():
            return {
                "status": "failed",
                "front": "",
                "middle": "",
                "rear": "",
                "method": method,
            }

        front, middle, rear = _split_text(text)
        return {
            "status": "success",
            "front": front,
            "middle": middle,
            "rear": rear,
            "method": method,
        }
    except Exception:
        return {
            "status": "failed",
            "front": "",
            "middle": "",
            "rear": "",
            "method": ext.lstrip(".") or "unknown",
        }
